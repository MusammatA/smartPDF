from __future__ import annotations

import cgi
import csv
import html
import io
import json
import mimetypes
import os
import re
import subprocess
import tempfile
from email import policy
from email.parser import BytesParser
from http import HTTPStatus
from http.server import SimpleHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path
from typing import Callable
from xml.etree import ElementTree as ET

from docx import Document
from docx.shared import Inches
from lxml import html as lxml_html
from openpyxl import Workbook, load_workbook
from PIL import Image, ImageOps
from pypdf import PdfReader
from reportlab.lib import colors
from reportlab.lib.pagesizes import LETTER, landscape
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image as ReportLabImage
from reportlab.platypus import Paragraph, Preformatted, SimpleDocTemplate, Spacer, Table, TableStyle

BASE_DIR = Path(__file__).resolve().parent
HOST = "127.0.0.1"
PORT = 8000
MAX_UPLOAD_BYTES = 50 * 1024 * 1024


class ConversionError(Exception):
    pass


SOURCE_EXTENSIONS = {
    "PDF": [".pdf"],
    "DOCX": [".docx"],
    "TXT": [".txt"],
    "HTML": [".html", ".htm"],
    "EPUB": [".epub"],
    "ODT": [".odt"],
    "RTF": [".rtf"],
    "MD": [".md", ".markdown"],
    "XLSX": [".xlsx"],
    "XLS": [".xls"],
    "ODS": [".ods"],
    "CSV": [".csv"],
    "TSV": [".tsv"],
    "JSON": [".json"],
    "XML": [".xml"],
    "PPTX": [".pptx"],
    "ODP": [".odp"],
    "PNG": [".png"],
    "JPG": [".jpg", ".jpeg"],
    "WEBP": [".webp"],
    "GIF": [".gif"],
    "BMP": [".bmp"],
    "TIFF": [".tif", ".tiff"],
    "ICO": [".ico"],
    "AVIF": [".avif"],
    "HEIC": [".heic", ".heif"],
    "SVG": [".svg"],
    "AI": [".ai"],
    "EPS": [".eps"],
    "DXF": [".dxf"],
    "Image": [
        ".png",
        ".jpg",
        ".jpeg",
        ".webp",
        ".gif",
        ".bmp",
        ".tif",
        ".tiff",
        ".ico",
        ".avif",
        ".heic",
        ".heif",
    ],
    "MP3": [".mp3"],
    "WAV": [".wav"],
    "AAC": [".aac"],
    "FLAC": [".flac"],
    "OGG": [".ogg"],
    "M4A": [".m4a"],
    "AIFF": [".aiff", ".aif"],
    "MP4": [".mp4"],
    "MOV": [".mov"],
    "AVI": [".avi"],
    "MKV": [".mkv"],
    "WEBM": [".webm"],
    "WMV": [".wmv"],
    "FLV": [".flv"],
    "SRT": [".srt"],
    "VTT": [".vtt"],
    "ASS": [".ass"],
    "ZIP": [".zip"],
    "RAR": [".rar"],
    "7Z": [".7z"],
    "TAR": [".tar"],
    "GZ": [".gz"],
    "XZ": [".xz"],
    "ISO": [".iso"],
    "MOBI": [".mobi"],
    "AZW3": [".azw3"],
    "FB2": [".fb2"],
    "TTF": [".ttf"],
    "OTF": [".otf"],
    "WOFF": [".woff"],
    "WOFF2": [".woff2"],
    "DWG": [".dwg"],
    "STL": [".stl"],
    "STEP": [".step", ".stp"],
    "OBJ": [".obj"],
    "IGES": [".iges", ".igs"],
    "FBX": [".fbx"],
    "GLTF": [".gltf"],
    "GLB": [".glb"],
    "PLY": [".ply"],
    "USDZ": [".usdz"],
    "EML": [".eml"],
    "MSG": [".msg"],
    "ICS": [".ics"],
    "GPX": [".gpx"],
    "KML": [".kml"],
    "GeoJSON": [".geojson", ".json"],
    "Text": [".txt", ".md", ".json", ".xml", ".csv", ".tsv", ".html", ".htm"],
    "Video": [".mp4", ".mov", ".avi", ".mkv", ".webm", ".wmv", ".flv"],
    "Audio": [".mp3", ".wav", ".aac", ".flac", ".ogg", ".m4a", ".aiff", ".aif"],
}

TARGET_EXTENSIONS = {
    "PDF": ".pdf",
    "DOCX": ".docx",
    "TXT": ".txt",
    "HTML": ".html",
    "EPUB": ".epub",
    "ODT": ".odt",
    "RTF": ".rtf",
    "MD": ".md",
    "XLSX": ".xlsx",
    "CSV": ".csv",
    "TSV": ".tsv",
    "JSON": ".json",
    "XML": ".xml",
    "PPTX": ".pptx",
    "ODP": ".odp",
    "PNG": ".png",
    "JPG": ".jpg",
    "WEBP": ".webp",
    "GIF": ".gif",
    "BMP": ".bmp",
    "TIFF": ".tiff",
    "ICO": ".ico",
    "AVIF": ".avif",
    "HEIC": ".heic",
    "SVG": ".svg",
    "AI": ".ai",
    "EPS": ".eps",
    "DXF": ".dxf",
    "MP3": ".mp3",
    "WAV": ".wav",
    "AAC": ".aac",
    "FLAC": ".flac",
    "OGG": ".ogg",
    "M4A": ".m4a",
    "AIFF": ".aiff",
    "MP4": ".mp4",
    "MOV": ".mov",
    "AVI": ".avi",
    "MKV": ".mkv",
    "WEBM": ".webm",
    "WMV": ".wmv",
    "FLV": ".flv",
    "SRT": ".srt",
    "VTT": ".vtt",
    "ASS": ".ass",
    "ZIP": ".zip",
    "RAR": ".rar",
    "7Z": ".7z",
    "TAR": ".tar",
    "GZ": ".gz",
    "XZ": ".xz",
    "ISO": ".iso",
    "MOBI": ".mobi",
    "AZW3": ".azw3",
    "FB2": ".fb2",
    "TTF": ".ttf",
    "OTF": ".otf",
    "WOFF": ".woff",
    "WOFF2": ".woff2",
    "DWG": ".dwg",
    "STL": ".stl",
    "STEP": ".step",
    "OBJ": ".obj",
    "IGES": ".iges",
    "FBX": ".fbx",
    "GLTF": ".gltf",
    "GLB": ".glb",
    "PLY": ".ply",
    "USDZ": ".usdz",
    "EML": ".eml",
    "MSG": ".msg",
    "ICS": ".ics",
    "GPX": ".gpx",
    "KML": ".kml",
    "GeoJSON": ".geojson",
}

CATALOG_INPUT = {
    "Documents": [
        "PDF ↔ DOCX",
        "PDF ↔ TXT",
        "PDF ↔ HTML",
        "PDF ↔ EPUB",
        "DOCX ↔ ODT",
        "DOCX ↔ RTF",
        "DOCX ↔ TXT",
        "DOCX ↔ HTML",
        "DOCX ↔ MD",
        "ODT ↔ PDF",
        "RTF ↔ TXT",
    ],
    "Spreadsheets": [
        "XLSX ↔ CSV",
        "XLSX ↔ ODS",
        "XLS ↔ XLSX",
        "XLSX ↔ PDF",
        "XLSX ↔ HTML",
        "CSV ↔ TSV",
        "CSV ↔ JSON",
    ],
    "Presentations": [
        "PPTX ↔ PDF",
        "PPTX ↔ ODP",
        "PPTX ↔ Images",
        "PPTX ↔ HTML",
    ],
    "Images (Raster)": [
        "PNG ↔ JPG",
        "PNG ↔ WEBP",
        "PNG ↔ GIF",
        "PNG ↔ BMP",
        "PNG ↔ TIFF",
        "PNG ↔ ICO",
        "PNG ↔ AVIF",
        "PNG ↔ HEIC",
        "JPG ↔ WEBP",
        "JPG ↔ TIFF",
        "GIF ↔ WEBP",
        "BMP ↔ TIFF",
    ],
    "Images (Vector)": [
        "SVG ↔ AI",
        "SVG ↔ EPS",
        "SVG ↔ PDF",
        "SVG ↔ DXF",
    ],
    "Raster ↔ Vector": [
        "PNG ↔ SVG",
        "JPG ↔ SVG",
        "BMP ↔ SVG",
    ],
    "Image Documents": [
        "Image ↔ PDF",
        "Image ↔ DOCX",
        "Image ↔ TXT",
        "Image ↔ HTML",
    ],
    "Audio": [
        "MP3 ↔ WAV",
        "MP3 ↔ AAC",
        "MP3 ↔ FLAC",
        "MP3 ↔ OGG",
        "MP3 ↔ M4A",
        "WAV ↔ FLAC",
        "WAV ↔ AIFF",
        "WAV ↔ AAC",
        "AAC ↔ M4A",
        "FLAC ↔ ALAC",
    ],
    "Video": [
        "MP4 ↔ MOV",
        "MP4 ↔ AVI",
        "MP4 ↔ MKV",
        "MP4 ↔ WEBM",
        "MP4 ↔ WMV",
        "MP4 ↔ FLV",
        "MP4 ↔ GIF",
        "MOV ↔ AVI",
        "MOV ↔ MKV",
        "AVI ↔ MKV",
        "MKV ↔ WEBM",
        "WEBM ↔ AVI",
    ],
    "Video Frames": [
        "Video ↔ GIF",
        "Video ↔ Images",
        "Images ↔ Video",
    ],
    "Subtitles": [
        "SRT ↔ VTT",
        "SRT ↔ ASS",
        "SRT ↔ TXT",
        "VTT ↔ ASS",
    ],
    "Archives": [
        "ZIP ↔ RAR",
        "ZIP ↔ 7Z",
        "ZIP ↔ TAR",
        "TAR ↔ GZ",
        "TAR ↔ XZ",
        "ISO ↔ ZIP",
    ],
    "eBooks": [
        "EPUB ↔ PDF",
        "EPUB ↔ MOBI",
        "EPUB ↔ AZW3",
        "EPUB ↔ HTML",
        "EPUB ↔ TXT",
        "FB2 ↔ EPUB",
    ],
    "Code & Markup": [
        "JSON ↔ XML",
        "JSON ↔ CSV",
        "XML ↔ CSV",
        "XML ↔ HTML",
        "HTML ↔ PDF",
        "HTML ↔ DOCX",
        "HTML ↔ MD",
        "MD ↔ HTML",
        "MD ↔ PDF",
        "MD ↔ DOCX",
        "SQL ↔ CSV",
    ],
    "Fonts": [
        "TTF ↔ OTF",
        "TTF ↔ WOFF",
        "TTF ↔ WOFF2",
        "OTF ↔ WOFF",
        "WOFF ↔ WOFF2",
    ],
    "CAD": [
        "DWG ↔ DXF",
        "DXF ↔ SVG",
        "STL ↔ STEP",
        "STL ↔ OBJ",
        "STEP ↔ IGES",
    ],
    "3D Models": [
        "OBJ ↔ FBX",
        "OBJ ↔ STL",
        "OBJ ↔ GLTF",
        "GLB ↔ FBX",
        "GLB ↔ STL",
        "PLY ↔ OBJ",
        "USDZ ↔ GLB",
    ],
    "Email": [
        "EML ↔ PDF",
        "MSG ↔ PDF",
        "EML ↔ HTML",
    ],
    "Calendars": [
        "ICS ↔ CSV",
        "ICS ↔ JSON",
    ],
    "GPS": [
        "GPX ↔ KML",
        "GPX ↔ GeoJSON",
        "KML ↔ GeoJSON",
    ],
    "OCR": [
        "PDF → TXT",
        "PDF → DOCX",
        "PDF → HTML",
        "Image → TXT",
        "Image → DOCX",
        "Image → PDF",
    ],
    "Speech": [
        "Audio ↔ Text",
        "Video ↔ Text",
    ],
    "AI": [
        "Text → Speech",
        "Text → Image",
        "Image → Text",
        "Image → HTML",
        "Image → React",
        "Sketch → SVG",
        "Sketch → PNG",
        "Table Image → XLSX",
        "Diagram Image → Mermaid",
        "Diagram Image → PlantUML",
        "Math Image → LaTeX",
        "Code Image → Source Code",
        "Handwriting → Text",
        "Black & White → Color",
        "Low Resolution → High Resolution",
    ],
}

STATUS_NOTES = {
    "available": "Runs locally in this app right now.",
    "partial": "Runs locally, but layout-heavy formats may export as simplified text or structure.",
    "unavailable": "Listed in the catalog, but this route needs extra local engines such as OCR, FFmpeg, Office renderers, or format-specific converters.",
}


def slugify(value: str) -> str:
    return re.sub(r"[^a-z0-9]+", "-", value.lower()).strip("-")


def safe_stem(filename: str) -> str:
    stem = Path(filename or "converted").stem
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "-", stem).strip("-._")
    return cleaned or "converted"


def extension_for_format(fmt: str) -> str:
    return TARGET_EXTENSIONS.get(fmt, f".{slugify(fmt)}")


def detect_suffix(value: str) -> str:
    return Path(value).suffix.lower()


def ensure_supported_input(upload_name: str, source_format: str) -> None:
    allowed = SOURCE_EXTENSIONS.get(source_format, [])
    if not allowed:
        return
    suffix = detect_suffix(upload_name)
    if suffix not in allowed:
        raise ConversionError(
            f"Expected a {source_format} upload. Allowed extensions: {', '.join(allowed)}"
        )


def read_text_file(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="latin-1")


def html_escape_paragraphs(text: str) -> str:
    pieces = []
    for block in re.split(r"\n\s*\n", text.strip()):
        if block.strip():
            pieces.append(f"<p>{html.escape(block).replace(chr(10), '<br>')}</p>")
    return "\n".join(pieces) or "<p></p>"


def markdown_inline(text: str) -> str:
    escaped = html.escape(text)
    escaped = re.sub(r"`([^`]+)`", r"<code>\1</code>", escaped)
    escaped = re.sub(r"\*\*([^*]+)\*\*", r"<strong>\1</strong>", escaped)
    escaped = re.sub(r"\*([^*]+)\*", r"<em>\1</em>", escaped)
    escaped = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r'<a href="\2">\1</a>', escaped)
    return escaped


def markdown_to_html(markdown_text: str) -> str:
    lines = markdown_text.splitlines()
    parts: list[str] = []
    in_list = False
    in_code = False
    code_lines: list[str] = []

    def close_list() -> None:
        nonlocal in_list
        if in_list:
            parts.append("</ul>")
            in_list = False

    for line in lines:
        stripped = line.rstrip()
        if stripped.startswith("```"):
            close_list()
            if in_code:
                parts.append("<pre><code>" + "\n".join(code_lines) + "</code></pre>")
                code_lines = []
                in_code = False
            else:
                in_code = True
            continue
        if in_code:
            code_lines.append(html.escape(line))
            continue
        if not stripped:
            close_list()
            continue
        heading = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if heading:
            close_list()
            level = len(heading.group(1))
            parts.append(f"<h{level}>{markdown_inline(heading.group(2))}</h{level}>")
            continue
        bullet = re.match(r"^[-*]\s+(.*)$", stripped)
        if bullet:
            if not in_list:
                parts.append("<ul>")
                in_list = True
            parts.append(f"<li>{markdown_inline(bullet.group(1))}</li>")
            continue
        close_list()
        parts.append(f"<p>{markdown_inline(stripped)}</p>")

    close_list()
    if in_code:
        parts.append("<pre><code>" + "\n".join(code_lines) + "</code></pre>")

    body = "\n".join(parts) or "<p></p>"
    return (
        "<!doctype html><html><head><meta charset='utf-8'>"
        "<title>Markdown Export</title></head><body>"
        f"{body}</body></html>"
    )


def html_file_to_tree(path: Path):
    text = read_text_file(path)
    return lxml_html.fromstring(text)


def html_to_plain_text(path: Path) -> str:
    tree = html_file_to_tree(path)
    text = tree.text_content()
    cleaned = re.sub(r"\n{3,}", "\n\n", text)
    return cleaned.strip()


def html_to_markdown(path: Path) -> str:
    tree = html_file_to_tree(path)
    body = tree.find("body") or tree
    lines: list[str] = []

    def render_element(el) -> None:
        tag = (el.tag or "").lower() if isinstance(el.tag, str) else ""
        text = " ".join(el.text_content().split())
        if tag in {"h1", "h2", "h3", "h4", "h5", "h6"} and text:
            level = int(tag[1])
            lines.append(f"{'#' * level} {text}")
            lines.append("")
            return
        if tag in {"p", "div"} and text:
            lines.append(text)
            lines.append("")
            return
        if tag == "li" and text:
            lines.append(f"- {text}")
            return
        if tag == "pre":
            lines.append("```")
            lines.append(el.text_content().strip())
            lines.append("```")
            lines.append("")
            return
        for child in el:
            render_element(child)

    for child in body:
        render_element(child)

    output = "\n".join(lines).strip()
    return output or body.text_content().strip()


def build_story_from_text(text: str, title: str | None = None):
    styles = getSampleStyleSheet()
    story = []
    if title:
        story.append(Paragraph(html.escape(title), styles["Title"]))
        story.append(Spacer(1, 0.2 * inch))
    for block in re.split(r"\n\s*\n", text.strip()):
        block = block.strip()
        if block:
            if "\n" in block:
                story.append(Preformatted(html.escape(block), styles["Code"]))
            else:
                story.append(Paragraph(html.escape(block), styles["BodyText"]))
            story.append(Spacer(1, 0.12 * inch))
    if not story:
        story.append(Paragraph("Empty document", styles["BodyText"]))
    return story


def write_text_pdf(text: str, out_path: Path, title: str | None = None) -> None:
    document = SimpleDocTemplate(str(out_path), pagesize=LETTER)
    story = build_story_from_text(text, title=title)
    document.build(story)


def text_to_docx(text: str, out_path: Path, heading: str | None = None) -> None:
    document = Document()
    if heading:
        document.add_heading(heading, level=1)
    for block in re.split(r"\n\s*\n", text.strip()):
        block = block.strip()
        if block:
            document.add_paragraph(block)
    if not document.paragraphs:
        document.add_paragraph("")
    document.save(out_path)


def docx_to_text(path: Path) -> str:
    document = Document(path)
    chunks: list[str] = []
    for paragraph in document.paragraphs:
        value = paragraph.text.strip()
        if value:
            chunks.append(value)
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                chunks.append(" | ".join(values))
    return "\n\n".join(chunks).strip()


def docx_to_html(path: Path) -> str:
    document = Document(path)
    parts = [
        "<!doctype html><html><head><meta charset='utf-8'>",
        "<title>DOCX Export</title></head><body>",
    ]
    for paragraph in document.paragraphs:
        text = html.escape(paragraph.text.strip())
        if not text:
            continue
        style_name = paragraph.style.name.lower() if paragraph.style and paragraph.style.name else ""
        if style_name.startswith("heading"):
            level_match = re.search(r"(\d+)", style_name)
            level = min(int(level_match.group(1)), 6) if level_match else 2
            parts.append(f"<h{level}>{text}</h{level}>")
        else:
            parts.append(f"<p>{text}</p>")
    for table in document.tables:
        parts.append("<table border='1' cellspacing='0' cellpadding='6'>")
        for row in table.rows:
            parts.append("<tr>")
            for cell in row.cells:
                parts.append(f"<td>{html.escape(cell.text.strip())}</td>")
            parts.append("</tr>")
        parts.append("</table>")
    parts.append("</body></html>")
    return "".join(parts)


def docx_to_markdown(path: Path) -> str:
    document = Document(path)
    lines: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = paragraph.style.name.lower() if paragraph.style and paragraph.style.name else ""
        if style_name.startswith("heading"):
            level_match = re.search(r"(\d+)", style_name)
            level = min(int(level_match.group(1)), 6) if level_match else 2
            lines.append(f"{'#' * level} {text}")
        else:
            lines.append(text)
        lines.append("")
    return "\n".join(lines).strip()


def pdf_to_text(path: Path) -> str:
    reader = PdfReader(str(path))
    parts = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            parts.append(text.strip())
    return "\n\n".join(parts).strip()


def write_html_file(content: str, out_path: Path) -> None:
    out_path.write_text(content, encoding="utf-8")


def run_textutil(src_path: Path, dest_path: Path, fmt: str) -> None:
    try:
        subprocess.run(
            ["/usr/bin/textutil", "-convert", fmt, str(src_path), "-output", str(dest_path)],
            check=True,
            capture_output=True,
            text=True,
        )
    except subprocess.CalledProcessError as exc:
        message = exc.stderr.strip() or exc.stdout.strip() or "textutil failed"
        raise ConversionError(message) from exc


def normalize_image_for_save(image: Image.Image, target_format: str) -> Image.Image:
    fmt = target_format.upper()
    if fmt in {"JPG", "JPEG", "AVIF", "BMP"}:
        if image.mode in {"RGBA", "LA"} or ("transparency" in image.info):
            background = Image.new("RGB", image.size, (255, 255, 255))
            alpha = image.convert("RGBA")
            background.paste(alpha, mask=alpha.split()[-1])
            return background
        return image.convert("RGB")
    if fmt in {"PNG", "WEBP", "GIF", "TIFF", "ICO"} and image.mode not in {"RGB", "RGBA", "P"}:
        return image.convert("RGBA")
    return image


def convert_image_format(src_path: Path, out_path: Path, target_format: str) -> None:
    with Image.open(src_path) as image:
        image = ImageOps.exif_transpose(image)
        prepared = normalize_image_for_save(image, target_format)
        save_kwargs = {}
        pillow_format = "JPEG" if target_format == "JPG" else target_format
        if target_format == "ICO":
            save_kwargs["sizes"] = [(256, 256)]
        prepared.save(out_path, format=pillow_format, **save_kwargs)


def image_to_pdf(src_path: Path, out_path: Path) -> None:
    with Image.open(src_path) as image:
        image = ImageOps.exif_transpose(image).convert("RGB")
        image.save(out_path, "PDF", resolution=144.0)


def image_to_docx(src_path: Path, out_path: Path) -> None:
    with Image.open(src_path) as image:
        width, height = image.size
    document = Document()
    document.add_heading("Image Export", level=1)
    max_width_inches = 6.5
    ratio = height / width if width else 1
    document.add_picture(str(src_path), width=Inches(max_width_inches))
    document.add_paragraph(f"Original size: {width} x {height}px")
    document.add_paragraph(f"Aspect ratio: {ratio:.2f}")
    document.save(out_path)


def image_to_html(src_path: Path, out_path: Path) -> None:
    mime_type = mimetypes.guess_type(src_path.name)[0] or "application/octet-stream"
    data = src_path.read_bytes()
    encoded = base64_encode(data)
    with Image.open(src_path) as image:
        width, height = image.size
    html_doc = f"""<!doctype html>
<html>
<head>
  <meta charset="utf-8">
  <title>{html.escape(src_path.stem)} Export</title>
  <style>
    body {{
      margin: 0;
      min-height: 100vh;
      background: #050505;
      color: #f5f5f5;
      font-family: "Avenir Next", "Helvetica Neue", sans-serif;
      display: grid;
      place-items: center;
      padding: 32px;
    }}
    .card {{
      background: #111;
      border: 1px solid #2e2e2e;
      padding: 24px;
      max-width: 860px;
      width: 100%;
    }}
    img {{
      width: 100%;
      height: auto;
      display: block;
      border: 1px solid #2e2e2e;
      background: #000;
    }}
  </style>
</head>
<body>
  <div class="card">
    <h1>{html.escape(src_path.name)}</h1>
    <p>{width} x {height}px</p>
    <img src="data:{mime_type};base64,{encoded}" alt="{html.escape(src_path.name)}">
  </div>
</body>
</html>
"""
    out_path.write_text(html_doc, encoding="utf-8")


def base64_encode(data: bytes) -> str:
    import base64

    return base64.b64encode(data).decode("ascii")


def worksheet_to_rows(path: Path) -> list[list[str]]:
    workbook = load_workbook(path, data_only=True)
    sheet = workbook.active
    rows = []
    for row in sheet.iter_rows(values_only=True):
        rows.append(["" if value is None else str(value) for value in row])
    return rows


def rows_to_csv(rows: list[list[str]], out_path: Path, delimiter: str = ",") -> None:
    with out_path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle, delimiter=delimiter)
        writer.writerows(rows)


def xlsx_to_csv(src_path: Path, out_path: Path) -> None:
    rows_to_csv(worksheet_to_rows(src_path), out_path, delimiter=",")


def xlsx_to_html(src_path: Path, out_path: Path) -> None:
    rows = worksheet_to_rows(src_path)
    table_rows = []
    for index, row in enumerate(rows):
        cells = []
        cell_tag = "th" if index == 0 else "td"
        for value in row:
            cells.append(f"<{cell_tag}>{html.escape(value)}</{cell_tag}>")
        table_rows.append("<tr>" + "".join(cells) + "</tr>")
    html_doc = (
        "<!doctype html><html><head><meta charset='utf-8'>"
        "<title>Spreadsheet Export</title></head><body><table border='1' cellspacing='0' cellpadding='6'>"
        + "".join(table_rows)
        + "</table></body></html>"
    )
    out_path.write_text(html_doc, encoding="utf-8")


def xlsx_to_pdf(src_path: Path, out_path: Path) -> None:
    rows = worksheet_to_rows(src_path)
    if not rows:
        raise ConversionError("The spreadsheet is empty.")
    document = SimpleDocTemplate(str(out_path), pagesize=landscape(LETTER))
    styles = getSampleStyleSheet()
    story = [
        Paragraph(html.escape(src_path.stem), styles["Title"]),
        Spacer(1, 0.2 * inch),
    ]
    widths = [1.4 * inch] * max(len(row) for row in rows)
    table = Table(rows, colWidths=widths, repeatRows=1)
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.black),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#444444")),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.HexColor("#f5f5f5"), colors.white]),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]
        )
    )
    story.append(table)
    document.build(story)


def csv_to_rows(src_path: Path, delimiter: str = ",") -> list[list[str]]:
    with src_path.open("r", encoding="utf-8", newline="") as handle:
        return [row for row in csv.reader(handle, delimiter=delimiter)]


def csv_to_json(src_path: Path, out_path: Path, delimiter: str = ",") -> None:
    with src_path.open("r", encoding="utf-8", newline="") as handle:
        reader = csv.DictReader(handle, delimiter=delimiter)
        rows = list(reader)
    out_path.write_text(json.dumps(rows, indent=2, ensure_ascii=False), encoding="utf-8")


def json_to_csv(src_path: Path, out_path: Path) -> None:
    data = json.loads(read_text_file(src_path))
    if isinstance(data, dict):
        if "rows" in data and isinstance(data["rows"], list):
            data = data["rows"]
        else:
            data = [data]
    if not isinstance(data, list) or not data:
        raise ConversionError("JSON to CSV expects a non-empty object list.")
    if not all(isinstance(item, dict) for item in data):
        raise ConversionError("JSON to CSV expects a list of objects.")
    headers: list[str] = []
    for item in data:
        for key in item.keys():
            if key not in headers:
                headers.append(key)
    with out_path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=headers)
        writer.writeheader()
        writer.writerows(data)


def json_to_xml(src_path: Path, out_path: Path) -> None:
    data = json.loads(read_text_file(src_path))
    root = ET.Element("root")

    def add_value(parent: ET.Element, key: str, value):
        if isinstance(value, dict):
            node = ET.SubElement(parent, key)
            for child_key, child_value in value.items():
                add_value(node, child_key, child_value)
        elif isinstance(value, list):
            node = ET.SubElement(parent, key)
            for item in value:
                add_value(node, "item", item)
        else:
            node = ET.SubElement(parent, key)
            node.text = "" if value is None else str(value)

    if isinstance(data, dict):
        for key, value in data.items():
            add_value(root, key, value)
    elif isinstance(data, list):
        for item in data:
            add_value(root, "item", item)
    else:
        add_value(root, "value", data)

    tree = ET.ElementTree(root)
    tree.write(out_path, encoding="utf-8", xml_declaration=True)


def xml_element_to_dict(element: ET.Element):
    children = list(element)
    if not children:
        return element.text or ""
    grouped = {}
    for child in children:
        value = xml_element_to_dict(child)
        if child.tag in grouped:
            current = grouped[child.tag]
            if not isinstance(current, list):
                grouped[child.tag] = [current]
            grouped[child.tag].append(value)
        else:
            grouped[child.tag] = value
    return grouped


def xml_to_json(src_path: Path, out_path: Path) -> None:
    root = ET.parse(src_path).getroot()
    payload = {root.tag: xml_element_to_dict(root)}
    out_path.write_text(json.dumps(payload, indent=2, ensure_ascii=False), encoding="utf-8")


def xml_to_csv(src_path: Path, out_path: Path) -> None:
    root = ET.parse(src_path).getroot()
    children = list(root)
    rows: list[dict[str, str]] = []
    if children and all(list(child) for child in children):
        for child in children:
            row = {}
            for grandchild in child:
                row[grandchild.tag] = (grandchild.text or "").strip()
            rows.append(row)
    else:
        row = {}
        for child in children:
            row[child.tag] = (child.text or "").strip()
        if row:
            rows.append(row)
    if not rows:
        raise ConversionError("XML to CSV needs child elements to map into rows.")
    headers: list[str] = []
    for row in rows:
        for key in row.keys():
            if key not in headers:
                headers.append(key)
    with out_path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=headers)
        writer.writeheader()
        writer.writerows(rows)


def csv_to_xml(src_path: Path, out_path: Path) -> None:
    with src_path.open("r", encoding="utf-8", newline="") as handle:
        rows = list(csv.DictReader(handle))
    root = ET.Element("rows")
    for row in rows:
        item = ET.SubElement(root, "row")
        for key, value in row.items():
            cell = ET.SubElement(item, key or "column")
            cell.text = value or ""
    tree = ET.ElementTree(root)
    tree.write(out_path, encoding="utf-8", xml_declaration=True)


def delimited_to_delimited(src_path: Path, out_path: Path, source_delimiter: str, target_delimiter: str) -> None:
    rows = csv_to_rows(src_path, delimiter=source_delimiter)
    rows_to_csv(rows, out_path, delimiter=target_delimiter)


def csv_to_xlsx(src_path: Path, out_path: Path, delimiter: str = ",") -> None:
    rows = csv_to_rows(src_path, delimiter=delimiter)
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Sheet1"
    for row in rows:
        sheet.append(row)
    workbook.save(out_path)


def html_to_docx(src_path: Path, out_path: Path) -> None:
    tree = html_file_to_tree(src_path)
    body = tree.find("body") or tree
    document = Document()

    def add_paragraph(text: str, style: str | None = None):
        if not text.strip():
            return
        if style:
            document.add_paragraph(text, style=style)
        else:
            document.add_paragraph(text)

    for child in body:
        tag = (child.tag or "").lower() if isinstance(child.tag, str) else ""
        text = " ".join(child.text_content().split())
        if not text:
            continue
        if tag in {"h1", "h2", "h3", "h4", "h5", "h6"}:
            document.add_heading(text, level=min(int(tag[1]), 4))
        elif tag == "li":
            document.add_paragraph(text, style="List Bullet")
        else:
            add_paragraph(text)

    if not document.paragraphs:
        document.add_paragraph("")
    document.save(out_path)


def html_to_pdf(src_path: Path, out_path: Path) -> None:
    text = html_to_plain_text(src_path)
    title = Path(src_path).stem.replace("-", " ").title()
    write_text_pdf(text, out_path, title=title)


def md_to_docx(src_path: Path, out_path: Path) -> None:
    html_path = out_path.with_suffix(".html")
    html_path.write_text(markdown_to_html(read_text_file(src_path)), encoding="utf-8")
    html_to_docx(html_path, out_path)
    html_path.unlink(missing_ok=True)


def md_to_pdf(src_path: Path, out_path: Path) -> None:
    html_path = out_path.with_suffix(".html")
    html_path.write_text(markdown_to_html(read_text_file(src_path)), encoding="utf-8")
    html_to_pdf(html_path, out_path)
    html_path.unlink(missing_ok=True)


def pdf_to_html(src_path: Path, out_path: Path) -> None:
    text = pdf_to_text(src_path)
    html_doc = (
        "<!doctype html><html><head><meta charset='utf-8'>"
        "<title>PDF Text Export</title></head><body>"
        f"{html_escape_paragraphs(text)}</body></html>"
    )
    out_path.write_text(html_doc, encoding="utf-8")


def pdf_to_docx(src_path: Path, out_path: Path) -> None:
    text_to_docx(pdf_to_text(src_path), out_path, heading=src_path.stem)


def pdf_to_txt(src_path: Path, out_path: Path) -> None:
    out_path.write_text(pdf_to_text(src_path), encoding="utf-8")


def docx_to_pdf(src_path: Path, out_path: Path) -> None:
    text = docx_to_text(src_path)
    write_text_pdf(text, out_path, title=src_path.stem)


def html_to_md(src_path: Path, out_path: Path) -> None:
    out_path.write_text(html_to_markdown(src_path), encoding="utf-8")


def md_to_html_file(src_path: Path, out_path: Path) -> None:
    out_path.write_text(markdown_to_html(read_text_file(src_path)), encoding="utf-8")


def docx_to_md(src_path: Path, out_path: Path) -> None:
    out_path.write_text(docx_to_markdown(src_path), encoding="utf-8")


def docx_to_html_file(src_path: Path, out_path: Path) -> None:
    out_path.write_text(docx_to_html(src_path), encoding="utf-8")


def txt_to_pdf(src_path: Path, out_path: Path) -> None:
    write_text_pdf(read_text_file(src_path), out_path, title=src_path.stem)


def txt_to_html(src_path: Path, out_path: Path) -> None:
    html_doc = (
        "<!doctype html><html><head><meta charset='utf-8'>"
        "<title>Text Export</title></head><body>"
        f"{html_escape_paragraphs(read_text_file(src_path))}</body></html>"
    )
    out_path.write_text(html_doc, encoding="utf-8")


def txt_to_docx(src_path: Path, out_path: Path) -> None:
    text_to_docx(read_text_file(src_path), out_path, heading=src_path.stem)


def txt_to_rtf(src_path: Path, out_path: Path) -> None:
    run_textutil(src_path, out_path, "rtf")


def html_to_eml(src_path: Path, out_path: Path) -> None:
    body = read_text_file(src_path)
    eml = (
        "Subject: HTML Export\n"
        "MIME-Version: 1.0\n"
        "Content-Type: text/html; charset=UTF-8\n"
        "\n"
        f"{body}"
    )
    out_path.write_text(eml, encoding="utf-8")


def eml_to_html(src_path: Path, out_path: Path) -> None:
    message = BytesParser(policy=policy.default).parsebytes(src_path.read_bytes())
    subject = html.escape(message.get("subject", "Email Export"))
    body = ""
    if message.is_multipart():
        for part in message.walk():
            content_type = part.get_content_type()
            if content_type == "text/html":
                body = part.get_content()
                break
            if content_type == "text/plain" and not body:
                body = "<pre>" + html.escape(part.get_content()) + "</pre>"
    else:
        if message.get_content_type() == "text/html":
            body = message.get_content()
        else:
            body = "<pre>" + html.escape(message.get_content()) + "</pre>"
    if not body:
        body = "<p>No readable message body found.</p>"
    out_path.write_text(
        "<!doctype html><html><head><meta charset='utf-8'>"
        f"<title>{subject}</title></head><body><h1>{subject}</h1>{body}</body></html>",
        encoding="utf-8",
    )


def eml_to_pdf(src_path: Path, out_path: Path) -> None:
    temp_html = out_path.with_suffix(".html")
    eml_to_html(src_path, temp_html)
    html_to_pdf(temp_html, out_path)
    temp_html.unlink(missing_ok=True)


def pptx_to_outline(path: Path) -> list[dict[str, object]]:
    from pptx import Presentation

    presentation = Presentation(path)
    slides: list[dict[str, object]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_lines.append(shape.text.strip())
        slides.append({"number": index, "content": slide_lines})
    return slides


def pptx_to_html(src_path: Path, out_path: Path) -> None:
    slides = pptx_to_outline(src_path)
    sections = []
    for slide in slides:
        items = "".join(f"<li>{html.escape(line)}</li>" for line in slide["content"])
        sections.append(
            f"<section><h2>Slide {slide['number']}</h2><ul>{items or '<li>(empty slide)</li>'}</ul></section>"
        )
    doc = (
        "<!doctype html><html><head><meta charset='utf-8'>"
        f"<title>{html.escape(src_path.stem)}</title></head><body>"
        f"<h1>{html.escape(src_path.stem)}</h1>{''.join(sections)}</body></html>"
    )
    out_path.write_text(doc, encoding="utf-8")


def pptx_to_pdf(src_path: Path, out_path: Path) -> None:
    slides = pptx_to_outline(src_path)
    blocks = []
    for slide in slides:
        lines = "\n".join(slide["content"]) if slide["content"] else "(empty slide)"
        blocks.append(f"Slide {slide['number']}\n{lines}")
    write_text_pdf("\n\n".join(blocks), out_path, title=src_path.stem)


def html_to_pptx(src_path: Path, out_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    tree = html_file_to_tree(src_path)
    body = tree.find("body") or tree
    slides_content: list[tuple[str, list[str]]] = []
    current_title = "Slide 1"
    current_points: list[str] = []
    for child in body:
        tag = (child.tag or "").lower() if isinstance(child.tag, str) else ""
        text = " ".join(child.text_content().split())
        if not text:
            continue
        if tag in {"h1", "h2"}:
            if current_points or slides_content:
                slides_content.append((current_title, current_points or [""]))
            current_title = text
            current_points = []
        else:
            current_points.append(text)
    slides_content.append((current_title, current_points or [""]))

    presentation = Presentation()
    for title, points in slides_content:
        layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = title
        body_shape = slide.placeholders[1].text_frame
        body_shape.clear()
        for index, point in enumerate(points):
            paragraph = body_shape.paragraphs[0] if index == 0 else body_shape.add_paragraph()
            paragraph.text = point
    if len(presentation.slides) > len(slides_content):
        while len(presentation.slides) > len(slides_content):
            r_id = presentation.slides._sldIdLst[-1].rId
            presentation.part.drop_rel(r_id)
            del presentation.slides._sldIdLst[-1]
    presentation.save(out_path)


def ics_to_events(path: Path) -> list[dict[str, str]]:
    events: list[dict[str, str]] = []
    current: dict[str, str] | None = None
    for raw_line in read_text_file(path).splitlines():
        line = raw_line.strip()
        if line == "BEGIN:VEVENT":
            current = {}
        elif line == "END:VEVENT" and current is not None:
            events.append(current)
            current = None
        elif current is not None and ":" in line:
            key, value = line.split(":", 1)
            current[key.split(";")[0]] = value
    return events


def ics_to_json(src_path: Path, out_path: Path) -> None:
    out_path.write_text(json.dumps(ics_to_events(src_path), indent=2, ensure_ascii=False), encoding="utf-8")


def ics_to_csv(src_path: Path, out_path: Path) -> None:
    events = ics_to_events(src_path)
    if not events:
        raise ConversionError("No calendar events found in the ICS file.")
    headers: list[str] = []
    for event in events:
        for key in event.keys():
            if key not in headers:
                headers.append(key)
    with out_path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=headers)
        writer.writeheader()
        writer.writerows(events)


def gpx_to_geojson(src_path: Path, out_path: Path) -> None:
    root = ET.parse(src_path).getroot()
    namespace = {"gpx": root.tag.partition("}")[0].strip("{")} if "}" in root.tag else {}
    selector = ".//gpx:trkpt" if namespace else ".//trkpt"
    points = []
    for point in root.findall(selector, namespace):
        lat = point.attrib.get("lat")
        lon = point.attrib.get("lon")
        if lat and lon:
            points.append([float(lon), float(lat)])
    if not points:
        raise ConversionError("No GPX track points were found.")
    geojson = {
        "type": "FeatureCollection",
        "features": [
            {
                "type": "Feature",
                "geometry": {"type": "LineString", "coordinates": points},
                "properties": {},
            }
        ],
    }
    out_path.write_text(json.dumps(geojson, indent=2), encoding="utf-8")


def geojson_to_kml(src_path: Path, out_path: Path) -> None:
    payload = json.loads(read_text_file(src_path))
    coordinates: list[list[float]] = []
    for feature in payload.get("features", []):
        geometry = feature.get("geometry", {})
        if geometry.get("type") == "LineString":
            coordinates = geometry.get("coordinates", [])
            break
    if not coordinates:
        raise ConversionError("GeoJSON to KML expects a LineString feature.")
    joined = " ".join(f"{lon},{lat},0" for lon, lat in coordinates)
    kml = f"""<?xml version="1.0" encoding="UTF-8"?>
<kml xmlns="http://www.opengis.net/kml/2.2">
  <Document>
    <Placemark>
      <LineString>
        <coordinates>{joined}</coordinates>
      </LineString>
    </Placemark>
  </Document>
</kml>
"""
    out_path.write_text(kml, encoding="utf-8")


def gpx_to_kml(src_path: Path, out_path: Path) -> None:
    temp_geojson = out_path.with_suffix(".geojson")
    gpx_to_geojson(src_path, temp_geojson)
    geojson_to_kml(temp_geojson, out_path)
    temp_geojson.unlink(missing_ok=True)


def kml_to_geojson(src_path: Path, out_path: Path) -> None:
    root = ET.parse(src_path).getroot()
    namespace = {"kml": root.tag.partition("}")[0].strip("{")} if "}" in root.tag else {}
    selector = ".//kml:coordinates" if namespace else ".//coordinates"
    coordinates_element = root.find(selector, namespace)
    if coordinates_element is None or not (coordinates_element.text or "").strip():
        raise ConversionError("No KML coordinates were found.")
    coordinates = []
    for chunk in coordinates_element.text.strip().split():
        lon, lat, *_ = chunk.split(",")
        coordinates.append([float(lon), float(lat)])
    geojson = {
        "type": "FeatureCollection",
        "features": [
            {
                "type": "Feature",
                "geometry": {"type": "LineString", "coordinates": coordinates},
                "properties": {},
            }
        ],
    }
    out_path.write_text(json.dumps(geojson, indent=2), encoding="utf-8")


def unsupported_message(source: str, target: str) -> str:
    if {"MP3", "WAV", "AAC", "FLAC", "OGG", "M4A", "AIFF", "Audio"} & {source, target}:
        return "Audio and speech routes need media encoders or transcription engines that are not bundled here."
    if {"MP4", "MOV", "AVI", "MKV", "WEBM", "WMV", "FLV", "Video"} & {source, target}:
        return "Video routes need FFmpeg-style renderers to preserve frames, codecs, and timing."
    if source == "Image" and target == "TXT":
        return "Image to text needs OCR, which is not installed in this local bundle."
    if "SVG" in {source, target} or {"AI", "EPS", "DXF"} & {source, target}:
        return "Vector and CAD routes need dedicated graphic or CAD converters."
    if {"ZIP", "RAR", "7Z", "ISO"} & {source, target}:
        return "Archive routes need additional archive engines for formats beyond built-in ZIP/TAR support."
    if {"PPTX", "ODP", "Images"} & {source, target}:
        return "Presentation image and ODP rendering typically need an office renderer, which is not healthy in this environment."
    if {"EPUB", "MOBI", "AZW3", "FB2"} & {source, target}:
        return "eBook routes need dedicated eBook packaging and rendering logic beyond this local MVP."
    if {"TTF", "OTF", "WOFF", "WOFF2"} & {source, target}:
        return "Font conversions need font-specific tooling that is not bundled here."
    if {"DWG", "STL", "STEP", "OBJ", "IGES", "FBX", "GLTF", "GLB", "PLY", "USDZ"} & {source, target}:
        return "CAD and 3D routes need specialized geometry converters."
    if source == "Text" and target == "Speech":
        return "Text to speech needs a speech synthesis engine for downloadable audio output."
    if target in {"React", "Mermaid", "PlantUML", "LaTeX", "Source Code", "Color", "High Resolution"}:
        return "AI-assisted routes need model-backed inference engines, which are outside this local-only build."
    if source in {"Black & White", "Low Resolution", "Handwriting", "Sketch", "Diagram Image", "Math Image", "Code Image", "Table Image"}:
        return "This AI route needs specialized model inference to produce high-quality results."
    return STATUS_NOTES["unavailable"]


Converter = Callable[[Path, Path], None]


def make_textutil(fmt: str) -> Converter:
    return lambda src, dst: run_textutil(src, dst, fmt)


SUPPORTED_HANDLERS: dict[tuple[str, str], tuple[Converter, str]] = {
    ("PDF", "TXT"): (pdf_to_txt, "partial"),
    ("PDF", "HTML"): (pdf_to_html, "partial"),
    ("PDF", "DOCX"): (pdf_to_docx, "partial"),
    ("DOCX", "PDF"): (docx_to_pdf, "partial"),
    ("DOCX", "ODT"): (make_textutil("odt"), "available"),
    ("ODT", "DOCX"): (make_textutil("docx"), "available"),
    ("DOCX", "RTF"): (make_textutil("rtf"), "available"),
    ("RTF", "DOCX"): (make_textutil("docx"), "available"),
    ("DOCX", "TXT"): (make_textutil("txt"), "available"),
    ("TXT", "DOCX"): (txt_to_docx, "available"),
    ("DOCX", "HTML"): (docx_to_html_file, "partial"),
    ("HTML", "DOCX"): (html_to_docx, "partial"),
    ("DOCX", "MD"): (docx_to_md, "partial"),
    ("MD", "DOCX"): (md_to_docx, "partial"),
    ("ODT", "PDF"): (lambda src, dst: html_to_pdf(_temp_convert_textutil(src, "html"), dst), "partial"),
    ("RTF", "TXT"): (make_textutil("txt"), "available"),
    ("TXT", "RTF"): (txt_to_rtf, "available"),
    ("TXT", "PDF"): (txt_to_pdf, "available"),
    ("TXT", "HTML"): (txt_to_html, "available"),
    ("HTML", "PDF"): (html_to_pdf, "partial"),
    ("PDF", "HTML"): (pdf_to_html, "partial"),
    ("HTML", "MD"): (html_to_md, "partial"),
    ("MD", "HTML"): (md_to_html_file, "available"),
    ("MD", "PDF"): (md_to_pdf, "partial"),
    ("XLSX", "CSV"): (xlsx_to_csv, "available"),
    ("CSV", "XLSX"): (csv_to_xlsx, "available"),
    ("XLSX", "PDF"): (xlsx_to_pdf, "available"),
    ("PDF", "XLSX"): (lambda src, dst: (_ for _ in ()).throw(ConversionError(unsupported_message("PDF", "XLSX"))), "unavailable"),
    ("XLSX", "HTML"): (xlsx_to_html, "available"),
    ("HTML", "XLSX"): (lambda src, dst: (_ for _ in ()).throw(ConversionError(unsupported_message("HTML", "XLSX"))), "unavailable"),
    ("CSV", "TSV"): (lambda src, dst: delimited_to_delimited(src, dst, ",", "\t"), "available"),
    ("TSV", "CSV"): (lambda src, dst: delimited_to_delimited(src, dst, "\t", ","), "available"),
    ("CSV", "JSON"): (csv_to_json, "available"),
    ("JSON", "CSV"): (json_to_csv, "available"),
    ("JSON", "XML"): (json_to_xml, "available"),
    ("XML", "JSON"): (xml_to_json, "available"),
    ("XML", "CSV"): (xml_to_csv, "partial"),
    ("CSV", "XML"): (csv_to_xml, "available"),
    ("PPTX", "HTML"): (pptx_to_html, "partial"),
    ("HTML", "PPTX"): (html_to_pptx, "partial"),
    ("PPTX", "PDF"): (pptx_to_pdf, "partial"),
    ("PNG", "JPG"): (lambda src, dst: convert_image_format(src, dst, "JPG"), "available"),
    ("JPG", "PNG"): (lambda src, dst: convert_image_format(src, dst, "PNG"), "available"),
    ("PNG", "WEBP"): (lambda src, dst: convert_image_format(src, dst, "WEBP"), "available"),
    ("WEBP", "PNG"): (lambda src, dst: convert_image_format(src, dst, "PNG"), "available"),
    ("PNG", "GIF"): (lambda src, dst: convert_image_format(src, dst, "GIF"), "available"),
    ("GIF", "PNG"): (lambda src, dst: convert_image_format(src, dst, "PNG"), "available"),
    ("PNG", "BMP"): (lambda src, dst: convert_image_format(src, dst, "BMP"), "available"),
    ("BMP", "PNG"): (lambda src, dst: convert_image_format(src, dst, "PNG"), "available"),
    ("PNG", "TIFF"): (lambda src, dst: convert_image_format(src, dst, "TIFF"), "available"),
    ("TIFF", "PNG"): (lambda src, dst: convert_image_format(src, dst, "PNG"), "available"),
    ("PNG", "ICO"): (lambda src, dst: convert_image_format(src, dst, "ICO"), "available"),
    ("ICO", "PNG"): (lambda src, dst: convert_image_format(src, dst, "PNG"), "available"),
    ("PNG", "AVIF"): (lambda src, dst: convert_image_format(src, dst, "AVIF"), "available"),
    ("AVIF", "PNG"): (lambda src, dst: convert_image_format(src, dst, "PNG"), "available"),
    ("JPG", "WEBP"): (lambda src, dst: convert_image_format(src, dst, "WEBP"), "available"),
    ("WEBP", "JPG"): (lambda src, dst: convert_image_format(src, dst, "JPG"), "available"),
    ("JPG", "TIFF"): (lambda src, dst: convert_image_format(src, dst, "TIFF"), "available"),
    ("TIFF", "JPG"): (lambda src, dst: convert_image_format(src, dst, "JPG"), "available"),
    ("GIF", "WEBP"): (lambda src, dst: convert_image_format(src, dst, "WEBP"), "available"),
    ("WEBP", "GIF"): (lambda src, dst: convert_image_format(src, dst, "GIF"), "available"),
    ("BMP", "TIFF"): (lambda src, dst: convert_image_format(src, dst, "TIFF"), "available"),
    ("TIFF", "BMP"): (lambda src, dst: convert_image_format(src, dst, "BMP"), "available"),
    ("Image", "PDF"): (image_to_pdf, "available"),
    ("Image", "DOCX"): (image_to_docx, "available"),
    ("Image", "HTML"): (image_to_html, "available"),
    ("EML", "HTML"): (eml_to_html, "available"),
    ("HTML", "EML"): (html_to_eml, "available"),
    ("EML", "PDF"): (eml_to_pdf, "partial"),
    ("ICS", "JSON"): (ics_to_json, "available"),
    ("ICS", "CSV"): (ics_to_csv, "available"),
    ("GPX", "GeoJSON"): (gpx_to_geojson, "available"),
    ("GPX", "KML"): (gpx_to_kml, "available"),
    ("KML", "GeoJSON"): (kml_to_geojson, "available"),
    ("GeoJSON", "KML"): (geojson_to_kml, "available"),
}


def _temp_convert_textutil(src_path: Path, fmt: str) -> Path:
    temp_dir = Path(tempfile.mkdtemp(prefix="smartpdf-textutil-"))
    out_path = temp_dir / f"{src_path.stem}.{fmt}"
    run_textutil(src_path, out_path, fmt)
    return out_path


def get_handler(source: str, target: str) -> tuple[Converter | None, str]:
    if (source, target) in SUPPORTED_HANDLERS:
        handler, status = SUPPORTED_HANDLERS[(source, target)]
        if status != "unavailable":
            return handler, status
    return None, "unavailable"


ALLOWED_SOURCE_FORMATS = (
    "PDF",
    "DOCX",
    "TXT",
    "JPG",
    "PNG",
    "HEIC",
    "MP4",
    "MP3",
    "MOV",
    "XLSX",
    "CSV",
    "PPTX",
    "ZIP",
    "EPUB",
)


def build_catalog() -> list[dict[str, object]]:
    entries: list[dict[str, object]] = []
    seen: set[str] = set()
    for category, conversions in CATALOG_INPUT.items():
        for conversion in conversions:
            match = re.match(r"^(.*?)\s*(↔|→)\s*(.*?)$", conversion)
            if not match:
                continue
            source = match.group(1).strip()
            arrow = match.group(2)
            target = match.group(3).strip()
            pairs = [(source, target)]
            if arrow == "↔":
                pairs.append((target, source))
            for from_fmt, to_fmt in pairs:
                if from_fmt not in ALLOWED_SOURCE_FORMATS:
                    continue
                label = f"{from_fmt} → {to_fmt}"
                key = slugify(label)
                if key in seen:
                    continue
                seen.add(key)
                _, status = get_handler(from_fmt, to_fmt)
                note = STATUS_NOTES[status] if status != "unavailable" else unsupported_message(from_fmt, to_fmt)
                entries.append(
                    {
                        "key": key,
                        "label": label,
                        "source": from_fmt,
                        "target": to_fmt,
                        "category": category,
                        "status": status,
                        "note": note,
                        "accept": SOURCE_EXTENSIONS.get(from_fmt, []),
                    }
                )
    return entries


CATALOG = build_catalog()
CATALOG_BY_KEY = {entry["key"]: entry for entry in CATALOG}


def convert_file(upload_name: str, payload: bytes, key: str) -> tuple[Path, str]:
    entry = CATALOG_BY_KEY.get(key)
    if not entry:
        raise ConversionError("Unknown conversion selected.")
    source = str(entry["source"])
    target = str(entry["target"])
    handler, status = get_handler(source, target)
    if handler is None or status == "unavailable":
        raise ConversionError(unsupported_message(source, target))
    ensure_supported_input(upload_name, source)

    request_dir = Path(tempfile.mkdtemp(prefix="smartpdf-job-"))
    source_extension = detect_suffix(upload_name) or extension_for_format(source)
    input_path = request_dir / f"input{source_extension}"
    input_path.write_bytes(payload)
    output_path = request_dir / f"{safe_stem(upload_name)}{extension_for_format(target)}"
    try:
        handler(input_path, output_path)
    except ConversionError:
        raise
    except Exception as exc:
        raise ConversionError(str(exc)) from exc
    if not output_path.exists():
        raise ConversionError("The conversion finished without producing a file.")
    return output_path, f"{safe_stem(upload_name)}{extension_for_format(target)}"


class SmartPDFHandler(SimpleHTTPRequestHandler):
    def log_message(self, format: str, *args) -> None:
        print(f"[smartPDF] {self.address_string()} - {format % args}")

    def do_GET(self) -> None:
        if self.path == "/api/conversions":
            payload = json.dumps({"conversions": CATALOG}, ensure_ascii=False).encode("utf-8")
            self.send_response(HTTPStatus.OK)
            self.send_header("Content-Type", "application/json; charset=utf-8")
            self.send_header("Content-Length", str(len(payload)))
            self.end_headers()
            self.wfile.write(payload)
            return
        if self.path == "/api/health":
            payload = json.dumps({"ok": True}).encode("utf-8")
            self.send_response(HTTPStatus.OK)
            self.send_header("Content-Type", "application/json; charset=utf-8")
            self.send_header("Content-Length", str(len(payload)))
            self.end_headers()
            self.wfile.write(payload)
            return
        if self.path == "/":
            self.path = "/index.html"
        return super().do_GET()

    def do_POST(self) -> None:
        if self.path != "/api/convert":
            self.send_error(HTTPStatus.NOT_FOUND, "Unknown endpoint")
            return
        length = int(self.headers.get("Content-Length", "0"))
        if length > MAX_UPLOAD_BYTES:
            self.respond_json(
                {"error": f"Upload is too large. Limit is {MAX_UPLOAD_BYTES // (1024 * 1024)} MB."},
                status=HTTPStatus.REQUEST_ENTITY_TOO_LARGE,
            )
            return
        form = cgi.FieldStorage(
            fp=self.rfile,
            headers=self.headers,
            environ={
                "REQUEST_METHOD": "POST",
                "CONTENT_TYPE": self.headers.get("Content-Type", ""),
            },
            keep_blank_values=True,
        )
        if "file" not in form or "conversionKey" not in form:
            self.respond_json({"error": "Both file and conversionKey are required."}, status=HTTPStatus.BAD_REQUEST)
            return
        upload_field = form["file"]
        conversion_key = form["conversionKey"].value
        if not getattr(upload_field, "file", None) or not upload_field.filename:
            self.respond_json({"error": "Please upload a file before converting."}, status=HTTPStatus.BAD_REQUEST)
            return
        payload = upload_field.file.read()
        if not payload:
            self.respond_json({"error": "The uploaded file was empty."}, status=HTTPStatus.BAD_REQUEST)
            return
        try:
            output_path, download_name = convert_file(upload_field.filename, payload, conversion_key)
        except ConversionError as exc:
            self.respond_json({"error": str(exc)}, status=HTTPStatus.BAD_REQUEST)
            return

        result = output_path.read_bytes()
        mime_type = mimetypes.guess_type(download_name)[0] or "application/octet-stream"
        self.send_response(HTTPStatus.OK)
        self.send_header("Content-Type", mime_type)
        self.send_header("Content-Disposition", f'attachment; filename="{download_name}"')
        self.send_header("Content-Length", str(len(result)))
        self.end_headers()
        self.wfile.write(result)

    def respond_json(self, payload: dict[str, object], status: HTTPStatus = HTTPStatus.OK) -> None:
        data = json.dumps(payload, ensure_ascii=False).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(data)))
        self.end_headers()
        self.wfile.write(data)


def main() -> None:
    os.chdir(BASE_DIR)
    server = ThreadingHTTPServer((HOST, PORT), SmartPDFHandler)
    print(f"smartPDF running at http://{HOST}:{PORT}")
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print("\nShutting down smartPDF...")
    finally:
        server.server_close()


if __name__ == "__main__":
    main()
